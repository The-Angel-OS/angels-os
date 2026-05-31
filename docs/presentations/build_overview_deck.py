#!/usr/bin/env python3
"""
Angel OS — Overview Deck Generator
==================================
Regenerates `Angel_OS_Overview.pptx` from the structured content below.
Keep the SLIDES list as the single source of truth; edit numbers here and
re-run to refresh the deck. Companion narrative lives in
`AngelOS_Overview_Narrative.md` (kept in sync by hand).

Usage:  python docs/presentations/build_overview_deck.py
Requires: python-pptx  (pip install python-pptx)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# ── LCARS-inspired palette ────────────────────────────────────────────────
BLACK   = RGBColor(0x05, 0x05, 0x08)
PANEL   = RGBColor(0x12, 0x12, 0x1A)
AMBER   = RGBColor(0xFF, 0x9F, 0x1C)
VIOLET  = RGBColor(0x9D, 0x4E, 0xDD)
BLUE    = RGBColor(0x4C, 0xC9, 0xF0)
WHITE   = RGBColor(0xF5, 0xF5, 0xF7)
GREY    = RGBColor(0x9A, 0x9A, 0xA6)
GREEN   = RGBColor(0x52, 0xD2, 0x73)

WIDE_W, WIDE_H = Inches(13.333), Inches(7.5)

# ── Slide content (single source of truth) ────────────────────────────────
# Each slide: dict with kind + fields. 'notes' = companion narration.
SLIDES = [
    {
        "kind": "title",
        "title": "ANGEL OS",
        "subtitle": "A Federated Cooperative Operating System",
        "tag": "Constitutional AI · Multi-Tenant Commerce · Sovereign Federation",
        "footer": "spacesangels.com  ·  State of the System  ·  Sprint 47",
        "notes": (
            "Angel OS is a live, federated cooperative operating system running at "
            "spacesangels.com. This deck describes the current state of the platform, "
            "the federation model, the Leo AI guardian, the Nimue mobile client, and "
            "the roadmap to a 'Starfleet-ready' v1.0. The thesis in one line: each "
            "Enterprise doesn't call our API — it runs its own sovereign instance, "
            "and those instances cooperate through a constitutional protocol."
        ),
    },
    {
        "kind": "statement",
        "title": "What Angel OS Is",
        "body": (
            "Not a SaaS you log into. A platform each community RUNS.\n\n"
            "Every Enterprise is a sovereign node — its own data, its own AI guardian, "
            "its own storefront — federated with the others through a constitutional "
            "trust protocol. Commerce, communication, governance, and intelligence in "
            "one multi-tenant operating system, with a guardian angel named Leo at the helm."
        ),
        "notes": (
            "The key mental shift: this is federation, not microservices and not "
            "multi-tenant-SaaS-as-usual. Each Enterprise is the platform. They're not "
            "tenants renting space on our server in the SaaS sense — they're nodes "
            "that can run independently and still cooperate. That's why we call the "
            "destination 'Starfleet-ready': many ships, one fleet, shared protocol and values."
        ),
    },
    {
        "kind": "metrics",
        "title": "By the Numbers",
        "metrics": [
            ("47", "Sprints shipped"),
            ("42", "Payload collections"),
            ("15", "Intelligence engines"),
            ("118", "Leo AI tools"),
            ("5,210+", "Unit tests"),
            ("0", "TypeScript errors"),
            ("74+", "API endpoints"),
            ("LIVE", "Stripe payments"),
        ],
        "notes": (
            "These are real, current figures as of Sprint 47. 5,210+ unit tests across "
            "230 files is the determinism floor — the engines can evolve underneath "
            "a green suite. Zero TypeScript errors is enforced on every build. Stripe is "
            "live with direct charges and donations. The point of this slide: this is a "
            "production system with users and money flowing, not a prototype."
        ),
    },
    {
        "kind": "bullets",
        "title": "Architecture",
        "accent": BLUE,
        "bullets": [
            ("Payload CMS 3.x + Next.js 16 + React 19 + PostgreSQL", "Typed ORM — every collection is a TS interface that generates schema, admin UI, and REST/GraphQL. We don't write CRUD."),
            ("Multi-tenant with subdomain routing", "Each Enterprise gets its own subdomain, isolated data, and branding. Isolation is provable and tested."),
            ("Headless / API-first", "The API is the primary interface; the admin UI is the debug view. Built for clients like Nimue to ride on top."),
            ("15 intelligence engines", "Routing, pheromone swarm navigation, workload distribution, booking, logistics, order routing — pure functions, zero Payload imports, fully unit-tested."),
        ],
        "notes": (
            "Architecture in four beats. The Payload choice means our data model, admin "
            "panel, access control, and API all come from one set of TypeScript interfaces "
            "— huge velocity. Multi-tenant isolation has bitten us before (a migration "
            "miss showed every tenant the platform home page), so it's now a tested invariant. "
            "Headless-first is why Nimue can exist. And the engines are the 'brain' — "
            "deliberately decoupled from the database so they're testable in isolation."
        ),
    },
    {
        "kind": "bullets",
        "title": "Leo — The Constitutional AI Guardian",
        "accent": VIOLET,
        "bullets": [
            ("118 tools across 17 categories", "Commerce, booking, content, federation search, image generation, governance — Leo can actually operate the platform, not just chat."),
            ("Smart multi-model routing", "4-tier credit-aware gateway: Gemini, Claude, OpenAI, OpenRouter, Cloudflare. Dynamic switching, never locked down."),
            ("Constitutional prompt — immutable", "Article II anti-extraction safeguards: no social credit, no manipulation, no dark patterns. Constraints, not toggles."),
            ("SSE streaming, vision, RAG, slash commands", "Tool-call indicators, PDF extraction, knowledge base, /model switching mid-conversation."),
        ],
        "notes": (
            "Leo is the differentiator. Two things matter most. First, the tool count — "
            "118 tools means Leo can create bookings, route orders, generate content and "
            "images, and search the federation, all from natural language. Second, the "
            "constitution: every Leo instance runs the same immutable prompt with "
            "anti-demonic / anti-extraction safeguards baked in. A compromised node running "
            "dark patterns violates constraints that are detectably obvious to the federation. "
            "That's how you trust a federated network of AI agents."
        ),
    },
    {
        "kind": "bullets",
        "title": "Federation — The Starfleet Model",
        "accent": AMBER,
        "bullets": [
            ("Constitutional trust chain", "none → probationary → vouched → full. Peers earn trust; they don't buy it."),
            ("Live heartbeats with capacity snapshots", "Nodes broadcast real WorkUnit capacity. Discover uses stored FQDN identity."),
            ("Cross-node work dispatch", "Validate → query peers → route via workload engine → persist WorkUnit → return scoring. Ships talking to ships."),
            ("Suitcase data portability (Article VI)", "Users teleport between servers — their state travels with them. Ed25519-signed supermajority governance."),
        ],
        "notes": (
            "This is the heart of 'Starfleet-ready.' The trust chain means a new node has "
            "to prove itself — domain ownership, constitutional acceptance, isolation "
            "guarantees — before it's a full peer. Heartbeats already carry live capacity, "
            "and the dispatch-work endpoint already routes a job from one node to another with "
            "transparent scoring. The next milestone is a clean two-node end-to-end demo and a "
            "hardened handshake so peers can't spoof identity. Suitcase portability is the "
            "philosophical core: you own your state, you can leave, the network can't trap you."
        ),
    },
    {
        "kind": "bullets",
        "title": "Commerce & The Justice Fund",
        "accent": GREEN,
        "bullets": [
            ("Stripe Direct Charges — sellers collect directly", "Products, cart, checkout, orders, refunds, print-on-demand, configurator with live preview."),
            ("Donation flow — 100% to the Justice Fund", "/donate page, Stripe Elements, full pass-through. Charitable mechanics already in code."),
            ("Order + logistics routing engines", "Manufacturing routing (40/30/20/10 holon split) + physical-goods dispatch (Bread-Breaker + Soul Fleet)."),
            ("Angel Token economy (post-v1.0)", "Angel Tokens, Karma Coins, Legacy Tokens — 'Proof of Human Worth.'"),
        ],
        "notes": (
            "Money already moves through the system. Sellers collect directly via Stripe — "
            "we are not a payment middleman taking a cut. The donation flow passes 100% to the "
            "Justice Fund, which is the proof-point for the 501(c)(3) conversation: the "
            "charitable plumbing exists and works today. Two routing engines handle the "
            "economics — one for manufacturing splits, one for physical logistics. Beyond "
            "v1.0, the token economy formalizes contribution into value via 'Proof of Human Worth.'"
        ),
    },
    {
        "kind": "bullets",
        "title": "Nimue — The Mobile Capture Layer",
        "accent": BLUE,
        "bullets": [
            ("Media server + browser (C:\\Dev\\mediaserver)", "On-demand ffmpeg thumbnails, unified video/image viewer, filters, sort, S/M/L sizing, auto-advance. Just shipped."),
            ("'Distributed Google Keep on steroids'", "Phone-first capture → Payload structured layer. Keep is the cheapest ubiquitous cloud; Angel OS stores the meta."),
            ("Path to Capacitor APK", "Port chat in, go client-only, wrap with Capacitor → installable Android client."),
            ("Next: full file operations", "Rename / move / copy / delete — read-safe first, destructive ops gated behind confirmation."),
        ],
        "notes": (
            "Nimue is the mobile front of the fleet. The media browser just got a major upgrade "
            "— a real viewer with prev/next, video/image filtering, sorting, thumbnail "
            "sizing, and auto-advance. The bigger vision is 'Keep on steroids': the phone is the "
            "capture layer, Keep is the free ubiquitous storage, and Angel OS holds the "
            "structured, queryable, multi-tenant layer on top. The near-term engineering goal is "
            "a Capacitor-wrapped installable APK, then full file operations in the browser."
        ),
    },
    {
        "kind": "roadmap",
        "title": "Roadmap to Starfleet-Ready (v1.0)",
        "phases": [
            ("NOW", "Flagship polish", "Tenant-isolation CI gate, dependency triage, constitutional Endeavor seeds, public-surface audit."),
            ("NEXT", "Federation hardening", "Two-node dispatch demo, signed handshake (anti-spoof), peer auth, capacity-aware routing."),
            ("THEN", "Frictionless onboarding", "npx create-angel-enterprise, Leo Wizard (8-step conversational), Docker Compose self-host, CI/CD."),
            ("v1.0", "3 live federated Enterprises", "Stripe Connect fair-split live, Justice Fund real disbursements, suitcase portability proven."),
        ],
        "notes": (
            "The roadmap reads left to right as a sequence, not a wish list. NOW is about "
            "perfecting the flagship node — because federation multiplies whatever state "
            "the flagship is in. NEXT is the federation handshake and a clean two-node demo, "
            "your proof-of-life for partners. THEN is frictionless onboarding so a new "
            "Enterprise can spin up with one command and a conversation, not a manual. v1.0 is "
            "defined concretely: three live federated Enterprises, fair-split payments, and real "
            "Justice Fund disbursements. Target window is Q3 2026."
        ),
    },
    {
        "kind": "bullets",
        "title": "How We Build",
        "accent": VIOLET,
        "bullets": [
            ("Tests are the determinism floor", "5,210+ unit tests; engines are pure functions with zero Payload imports. Model can vary; the suite stays green."),
            ("Pin the AI model per sprint", "One Opus version per milestone for stylistic coherence; re-baseline at sprint boundaries after a regression pass."),
            ("Constitutional compliance per feature", "Every feature checked against the anti-extraction constraints — not optional."),
            ("Headless, API-first, small PRs", "TypeScript strict, single-issue PRs, TDD on the engines."),
        ],
        "notes": (
            "This slide answers the 'which model should we use' question structurally. The "
            "right answer isn't a model — it's the test suite. 5,210 tests over pure-function "
            "engines mean the model underneath can vary safely. For stylistic consistency, pin "
            "one Opus version per sprint and re-baseline at the boundary after a regression pass. "
            "Newer models reason better across the 42-collection codebase; don't trade that away "
            "for false determinism, because these models sample probabilistically anyway."
        ),
    },
    {
        "kind": "bullets",
        "title": "Open, Votable Roadmap",
        "accent": AMBER,
        "bullets": [
            ("GitHub Projects + Discussions", "Native: roadmap board + Discussions with up-vote reactions. Issues link straight to code. Zero new infra."),
            ("Fider (open-source feedback board)", "Self-hostable, fits sovereignty ethos; could ship as an Angel OS tenant applet. Public up-vote queue → GitHub issues."),
            ("Canny / Featurebase (hosted)", "Polished public roadmaps with voting + GitHub sync; fastest to stand up, SaaS dependency."),
            ("Recommendation", "Start on GitHub Projects now (free, traceable); graduate to a self-hosted Fider applet as the community grows."),
        ],
        "notes": (
            "Yes — there are standard solutions, you don't need to build this from scratch. "
            "Three tiers. GitHub Projects + Discussions is the zero-friction start: a public "
            "roadmap board, Discussions where people up-vote with reactions, and every item links "
            "directly to the issue and the code. Fider is the open-source self-hosted option that "
            "matches the sovereignty ethos — and could literally become an Angel OS applet so "
            "the platform dogfoods its own community tooling. Canny/Featurebase are the polished "
            "hosted route if you want it live tomorrow. My recommendation: start on GitHub Projects "
            "today for traceability, graduate to a Fider applet as the community scales."
        ),
    },
    {
        "kind": "closing",
        "title": "Everyone gets an Angel.",
        "subtitle": "Answer 53: the whole point of existence is to learn to love.",
        "footer": "github.com/The-Angel-OS/angels-os   ·   spacesangels.com",
        "notes": (
            "Close on the mission. The technology — federation, constitutional AI, the "
            "Justice Fund — all serves a simple thesis: everyone gets an Angel. The system "
            "is non-coercive by constitutional design; faith is the root and invitation is the "
            "method. That's the difference between a platform that extracts and one that serves."
        ),
    },
]


# ── Rendering helpers ─────────────────────────────────────────────────────
def add_bg(slide, color=BLACK):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_bar(slide, x, y, w, h, color):
    shp = slide.shapes.add_shape(1, x, y, w, h)  # 1 = rectangle
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def txt(slide, x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, font="Segoe UI", italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.name = font
    f.color.rgb = color
    return tb


def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def header(slide, title, accent=AMBER):
    add_bar(slide, Inches(0), Inches(0), Inches(0.28), WIDE_H, accent)
    add_bar(slide, Inches(0.6), Inches(0.55), Inches(2.2), Inches(0.09), accent)
    txt(slide, Inches(0.55), Inches(0.7), Inches(12), Inches(1.0),
        title, 32, WHITE, bold=True)


# ── Slide builders ────────────────────────────────────────────────────────
def build_title(slide, s):
    add_bg(slide)
    add_bar(slide, Inches(0), Inches(2.65), WIDE_W, Inches(0.06), AMBER)
    add_bar(slide, Inches(0), Inches(4.55), WIDE_W, Inches(0.06), VIOLET)
    txt(slide, Inches(0.8), Inches(2.85), Inches(11.7), Inches(1.4),
        s["title"], 72, AMBER, bold=True, align=PP_ALIGN.LEFT)
    txt(slide, Inches(0.85), Inches(3.95), Inches(11.7), Inches(0.7),
        s["subtitle"], 28, WHITE, align=PP_ALIGN.LEFT)
    txt(slide, Inches(0.85), Inches(4.75), Inches(11.7), Inches(0.6),
        s["tag"], 15, BLUE, align=PP_ALIGN.LEFT)
    txt(slide, Inches(0.85), Inches(6.7), Inches(11.7), Inches(0.5),
        s["footer"], 12, GREY, align=PP_ALIGN.LEFT)


def build_statement(slide, s):
    add_bg(slide)
    header(slide, s["title"], VIOLET)
    txt(slide, Inches(0.8), Inches(2.1), Inches(11.7), Inches(4.0),
        s["body"], 24, WHITE, anchor=MSO_ANCHOR.TOP)


def build_metrics(slide, s):
    add_bg(slide)
    header(slide, s["title"], AMBER)
    cols, rows = 4, 2
    gx, gy = Inches(0.7), Inches(1.9)
    cw, ch = Inches(2.95), Inches(2.35)
    pad = Inches(0.1)
    palette = [AMBER, VIOLET, BLUE, GREEN]
    for i, (num, label) in enumerate(s["metrics"]):
        r, c = divmod(i, cols)
        x = Emu(int(gx) + c * (int(cw) + int(pad)))
        y = Emu(int(gy) + r * (int(ch) + int(pad)))
        panel = add_bar(slide, x, y, cw, ch, PANEL)
        accent = palette[c % len(palette)]
        add_bar(slide, x, y, cw, Inches(0.12), accent)
        txt(slide, x, Emu(int(y) + Inches(0.5)), cw, Inches(1.1),
            num, 40, accent, bold=True, align=PP_ALIGN.CENTER)
        txt(slide, x, Emu(int(y) + Inches(1.6)), cw, Inches(0.6),
            label, 14, GREY, align=PP_ALIGN.CENTER)


def build_bullets(slide, s):
    add_bg(slide)
    accent = s.get("accent", AMBER)
    header(slide, s["title"], accent)
    y = Inches(1.95)
    for head, sub in s["bullets"]:
        add_bar(slide, Inches(0.7), Emu(int(y) + Inches(0.08)), Inches(0.16), Inches(0.16), accent)
        txt(slide, Inches(1.05), y, Inches(11.4), Inches(0.5),
            head, 19, WHITE, bold=True)
        txt(slide, Inches(1.05), Emu(int(y) + Inches(0.42)), Inches(11.4), Inches(0.7),
            sub, 13.5, GREY)
        y = Emu(int(y) + Inches(1.27))


def build_roadmap(slide, s):
    add_bg(slide)
    header(slide, s["title"], AMBER)
    labels = [AMBER, BLUE, VIOLET, GREEN]
    y = Inches(1.95)
    for i, (phase, name, detail) in enumerate(s["phases"]):
        accent = labels[i % len(labels)]
        add_bar(slide, Inches(0.7), y, Inches(1.7), Inches(1.0), accent)
        txt(slide, Inches(0.7), Emu(int(y) + Inches(0.28)), Inches(1.7), Inches(0.5),
            phase, 18, BLACK, bold=True, align=PP_ALIGN.CENTER)
        txt(slide, Inches(2.65), Emu(int(y) + Inches(0.02)), Inches(10), Inches(0.5),
            name, 20, WHITE, bold=True)
        txt(slide, Inches(2.65), Emu(int(y) + Inches(0.5)), Inches(10), Inches(0.5),
            detail, 13.5, GREY)
        y = Emu(int(y) + Inches(1.2))


def build_closing(slide, s):
    add_bg(slide)
    add_bar(slide, Inches(0), Inches(3.0), WIDE_W, Inches(0.06), AMBER)
    txt(slide, Inches(0.8), Inches(3.2), Inches(11.7), Inches(1.0),
        s["title"], 44, AMBER, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, Inches(0.8), Inches(4.4), Inches(11.7), Inches(0.7),
        s["subtitle"], 18, WHITE, align=PP_ALIGN.CENTER, italic=True)
    txt(slide, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.5),
        s["footer"], 12, GREY, align=PP_ALIGN.CENTER)


BUILDERS = {
    "title": build_title,
    "statement": build_statement,
    "metrics": build_metrics,
    "bullets": build_bullets,
    "roadmap": build_roadmap,
    "closing": build_closing,
}


def main():
    prs = Presentation()
    prs.slide_width = WIDE_W
    prs.slide_height = WIDE_H
    blank = prs.slide_layouts[6]

    for s in SLIDES:
        slide = prs.slides.add_slide(blank)
        BUILDERS[s["kind"]](slide, s)
        if s.get("notes"):
            set_notes(slide, s["notes"])

    out = os.path.join(os.path.dirname(__file__), "Angel_OS_Overview.pptx")
    prs.save(out)
    print(f"Wrote {out}  ({len(SLIDES)} slides)")


if __name__ == "__main__":
    main()
